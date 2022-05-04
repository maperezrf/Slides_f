from pptx import Presentation
from pptx.util import Cm,Pt
from f3 import F3
from f4_classifier import CLASSIFIER_F4
from f4 import F4
from datetime import datetime
from general import generate_structure

# A tener en cuenta: 
# TODO 1. Entre método y método solo un espacio de separación 
# TODO 2. cuando hay signo "=" dejar un espacio antes y después  
# TODO 3. cuando vayan "," que queden pegadas a la palabra anterior y con un espacio después 

fecha_corte = input("Ingrese la fecha de corte, formato -> AAMMDD: ") #[x] ingresar fecha por usuario AAAA-MM-DD

# Datos a modificar antes de la ejecución del código

fecha_riesgo_f3 = "2022-03-30" 

generate_structure(fecha_corte)

# Inicio
f3 = F3(fecha_riesgo_f3,fecha_corte)
path_f3 = f3.get_path() #[x] pasar a método get_path()

f4_classifier = CLASSIFIER_F4(fecha_corte)
f4_clasificada = f4_classifier.f4_clas_marc #TODO pasar a método get_f4_classified()

f4 = F4()
f4.iniciar(f4_clasificada,fecha_corte)
path_f4 = f4.path #TODO pasar a método get_path()
#[x] unificar con fecha de corte de fecha_corte_f3
f4_cd,f4_tienda,f4_dvd,f4_venta,reservado = f4_classifier.calculos()

seguimiento_fs = Presentation("input/plantilla_seg_fs.pptx")  # Leer presentacion 

def update_date(n_slide,fecha_corte): 
    slide = seguimiento_fs.slides[n_slide]
    tb_corte = slide.shapes.add_textbox(Cm(0),Cm(17.43),width=Cm(1), height=Cm(0.1))
    text_corte = tb_corte.text_frame.add_paragraph() 
    text_corte.text = f"Corte {fecha_corte}"
    text_corte.font.size = Pt(12)    

for i in  range(8,17):
    update_date(i,fecha_corte)

def slide_f3_costo():
    f3 = seguimiento_fs.slides[7]
    f3.shapes.add_picture(f"{path_f3}/{fecha_corte}_f3_abiertos_fecha_reserva.png",Cm(0.9),Cm(1.81),width=Cm(10.78), height=Cm(7.5))
    f3.shapes.add_picture(f"{path_f3}/{fecha_corte}_f3_Cerrado_producto_costo.png",Cm(13.52),Cm(1.37),width=Cm(9.16), height=Cm(7.94))
    f3.shapes.add_picture(f"{path_f3}/{fecha_corte}_f3_Cerrado_producto_costo.png",Cm(23.74),Cm(1.42),width=Cm(8.65), height=Cm(8.3))
    f3.shapes.add_picture(f"{path_f3}/{fecha_corte}_f3_abierto_sede.png",Cm(1.12),Cm(10.88),width=Cm(11.47), height=Cm(6.7))
    f3.shapes.add_picture(f"{path_f3}/{fecha_corte}_f3_tendencia_mkp.png",Cm(13.28),Cm(10.88),width=Cm(10.88), height=Cm(6.7))
    f3.shapes.add_picture(f"{path_f3}/{fecha_corte}_f3_Cerrado_mkp_costo.png",Cm(23.74),Cm(10.18),width=Cm(9.25), height=Cm(7.71))
    update_date(7, fecha_corte)

def slide_f4():
    f4 = seguimiento_fs.slides[8]
    f4.shapes.add_picture(f"{path_f4}/{fecha_corte}_f4_tendencia_creacion_f4_x_años.png", Cm(2.01),Cm(0.58),width = Cm(16.17), height = Cm(9.08)   )
    f4.shapes.add_picture(f"{path_f4}/{fecha_corte}_f4_grafica_f4_sem.png", Cm(19.55), Cm(0.57), width = Cm(13.29), height = Cm(9.6))
    f4.shapes.add_picture(f"{path_f4}/{fecha_corte}_f4_grafica_total_por_mes.png", Cm(14.03), Cm(9.91), width = Cm(7.62), height = Cm(7.62)   )
    f4.shapes.add_picture(f"{path_f4}/{fecha_corte}_f4_torta.png", Cm(20.85), Cm(10.69), width = Cm(9.82), height = Cm(7))
    tb_reserva = f4.shapes.add_textbox(Cm(1.21),Cm(15.28),width=Cm(1), height=Cm(0.1))
    tx_text_reserva = tb_reserva.text_frame.add_paragraph() 
    tx_text_reserva.text = f"Reservado {reservado}"
    tx_text_reserva.font.size= Pt(14)
    
    corchete = f4.shapes.add_textbox(Cm(5.03),Cm(14.05),width=Cm(1), height=Cm(0.1))
    tx_corchete = corchete.text_frame.add_paragraph() 
    tx_corchete.text = "{"
    tx_corchete.font.size= Pt(66)

    estados = f4.shapes.add_textbox(Cm(5.87),Cm(14.6),width=Cm(1), height=Cm(0.1))
    tx_estados = estados.text_frame.add_paragraph() 
    tx_estados.text = f"Tienda {f4_tienda}\nCD {f4_cd}\nDVD {f4_dvd}\nVenta empresa {f4_venta}"
    tx_estados.font.size= Pt(12)

    registrados = f4.shapes.add_textbox(Cm(1.24),Cm(8.8),width=Cm(0.1), height=Cm(0.1))
    tx_registrados = registrados.text_frame.add_paragraph() 
    tx_registrados.text = f"Estado Registrado $21M < 7 días No han pasado a contabilidad,\nya informados a responsables\nUsuarios de creación:\n    -Jhonatan Beltran Velasquez\n   -Alexis Fabian Mestizo Parra"
    tx_registrados.font.size= Pt(13)

    registrados_m = f4.shapes.add_textbox(Cm(1.24),Cm(12.08),width=Cm(0.1), height=Cm(0.1))
    tx_registrados_m = registrados_m.text_frame.add_paragraph() 
    tx_registrados_m.text = f"Estado Registrado $2M >7 días No han pasado a contabilidad,\nya informados a responsables\nUsuarios de creación:\n  Jair Armando Rocha Vargas"
    tx_registrados_m.font.size= Pt(13)       

def f4_pos_causa():
    f4_pos_causa = seguimiento_fs.slides[9]
    f4_pos_causa.shapes.add_picture(f"{path_f4}/{fecha_corte}_f4_clasificacion_posibles_causas_22.png", Cm(1.36),Cm(0.38),width = Cm(31.91), height = Cm(17.32))
    titulo = f4_pos_causa.shapes.add_textbox(Cm(0.4),Cm(-0.79),width=Cm(0.1), height=Cm(0.1))
    tx_titulo = titulo.text_frame.add_paragraph() 
    tx_titulo.text = "F4"
    tx_titulo.font.size= Pt(36)       

def f4_mot_sede():
    f4_mot_sede = seguimiento_fs.slides[10]
    f4_mot_sede.shapes.add_picture(f"{path_f4}/{fecha_corte}_f4_grafica_total.png", Cm(0.66),Cm(2.32),width = Cm(17.41), height = Cm(14.81))
    f4_mot_sede.shapes.add_picture(f"{path_f4}/{fecha_corte}_f4_grafica_total_por_mes.png", Cm(22.38), Cm(1.01), width = Cm(7.62), height = Cm(7.62))
    f4_mot_sede.shapes.add_picture(f"{path_f4}/{fecha_corte}_f4_grafica_total_por_local.png", Cm(18.55), Cm(8.52), width = Cm(14.18), height = Cm(8.94))

def f4_mes_moti():
    f4_mes_moti = seguimiento_fs.slides[11]
    f4_mes_moti.shapes.add_picture(f"{path_f4}/{fecha_corte}_f4_tienda_mes_motivo.png", Cm(0.34),Cm(2.88),width = Cm(16.03), height = Cm(14.02))
    f4_mes_moti.shapes.add_picture(f"{path_f4}/{fecha_corte}_f4s_cd_mm.png", Cm(19.5),Cm(0.61),width = Cm(11.59), height = Cm(8.27))
    f4_mes_moti.shapes.add_picture(f"{path_f4}/{fecha_corte}_f4_tienda_mes_motivo.png", Cm(19.13),Cm(9.19),width = Cm(11.72), height = Cm(8.38))

def f4_linea_año():
    f4_linea_año = seguimiento_fs.slides[12]
    f4_linea_año.shapes.add_picture(f"{path_f4}/{fecha_corte}_f4_torta_linea.png", Cm(1.61),Cm(2.26),width = Cm(13.18), height = Cm(15.06))
    f4_linea_año.shapes.add_picture(f"{path_f4}/{fecha_corte}_f4_linea_local.png", Cm(14.79),Cm(1.89),width = Cm(17.46), height = Cm(15.27))

def f4_linea_motivo():
    f4_linea_motivo = seguimiento_fs.slides[13]
    f4_linea_motivo.shapes.add_picture(f"{path_f4}/{fecha_corte}_f4_linea_motivo.png", Cm(0.21),Cm(2.39),width = Cm(16.47), height = Cm(13.19))
    f4_linea_motivo.shapes.add_picture(f"{path_f4}/{fecha_corte}_f4_grafica_linea_x_mes.png", Cm(15.56),Cm(1.61),width = Cm(17.66), height = Cm(15.41))

def f4_averias():
    f4_linea_motivo = seguimiento_fs.slides[14]
    f4_linea_motivo.shapes.add_picture(f"{path_f4}/{fecha_corte}_f4_grafica_averia_x_mes_y_marca.png", Cm(0),Cm(2.4),width = Cm(16.64), height = Cm(13.82))

def f4_calidad():
    f4_calidad = seguimiento_fs.slides[15]
    f4_calidad.shapes.add_picture(f"{path_f4}/{fecha_corte}_f4_marca_calidad.png", Cm(0),Cm(2.4),width = Cm(16.64), height = Cm(13.82))

def f4_panatallas_rotas():
    f4_panatallas_rotas = seguimiento_fs.slides[16]
    f4_panatallas_rotas.shapes.add_picture(f"{path_f4}/{fecha_corte}_f4_pantallas_rotas.png", Cm(0),Cm(4.1),width = Cm(19.73), height = Cm(12.33))
        
slide_f3_costo()    
slide_f4()
f4_pos_causa()
f4_mot_sede()
f4_mes_moti()
f4_linea_año()
f4_linea_motivo()
f4_averias()
f4_calidad()
f4_panatallas_rotas()
seguimiento_fs.save(f"output/{fecha_corte}seguimiento_fs.pptx")


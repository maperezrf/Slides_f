from pickletools import read_uint1
from sys import path_hooks
from pptx import Presentation
from pptx.util import Cm,Pt
from data import  var_global,var_main
from datetime import datetime
from pptx.dml.color import RGBColor
from data import var_f11
import calendar
from pptx.enum.shapes import MSO_SHAPE

class PPTPY():
    path_f3 = ''
    path_f4 = ''
    fecha_corte = ''
    f4_data = []

    def __init__(self, pf3, pf4, pf11, f4cals, fc, años) -> None:
        self.path_f3 = pf3
        self.path_f4 = pf4
        self.path_f11 = pf11
        self.fecha_corte = fc
        self.f4_data = f4cals
        self.años = años
        
        self.seguimiento_fs = Presentation(var_main['pah_plantilla'])  # Leer presentacion 
        self.update_date_ftont()
        for i in  range(2,33):
            self.update_date(i)
    
    def update_date_ftont(self): 
        slide = self.seguimiento_fs.slides[0]
        tb_corte = slide.shapes.add_textbox(Cm(26.48),Cm(12.34),width=Cm(1), height=Cm(0.2))
        text_corte = tb_corte.text_frame.add_paragraph() 
        text_corte.text = f"{datetime.strptime(self.fecha_corte, '%y%m%d').strftime('%d - %b - %Y')}"
        text_corte.font.size = Pt(25)
        text_corte.font.color.rgb = RGBColor(255, 255, 255)
        
    def update_date(self, n_slide): 
        slide = self.seguimiento_fs.slides[n_slide]
        tb_corte = slide.shapes.add_textbox(Cm(0),Cm(17.43),width=Cm(1), height=Cm(0.1))
        text_corte = tb_corte.text_frame.add_paragraph() 
        text_corte.text = f"Corte {datetime.strptime(self.fecha_corte, '%y%m%d').strftime('%Y-%m-%d')}"
        text_corte.font.size = Pt(12)    

    def slide_f11_empresa(self, slide, list_graf):
        f11_empresa = self.seguimiento_fs.slides[slide]
        f11_empresa.shapes.add_picture(f"{self.path_f11}/{self.fecha_corte}{list_graf[0]}", Cm(0), Cm(2.3), height = Cm(14.46)) 
        f11_empresa.shapes.add_picture(f"{self.path_f11}/{self.fecha_corte}{list_graf[1]}", Cm(15.86), Cm(1.73), height = Cm(5.75), width = Cm(8.53)) 
        f11_empresa.shapes.add_picture(f"{self.path_f11}/{self.fecha_corte}{list_graf[2]}", Cm(24.9), Cm(1.73), height = Cm(5.75), width = Cm(8.53)) 
        f11_empresa.shapes.add_picture(f"{self.path_f11}/{self.fecha_corte}{list_graf[3]}", Cm(15.86), Cm(10.18), height = Cm(5.75), width = Cm(8.53)) 
        f11_empresa.shapes.add_picture(f"{self.path_f11}/{self.fecha_corte}{list_graf[4]}", Cm(24.89), Cm(10.18),  height = Cm(5.75), width = Cm(8.53)) 
        shapes = f11_empresa.shapes
        shape = shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Cm(2.15), Cm(2.75), Cm(10.22), Cm(10.23)
        )
        shape.fill.background()
        line = shape.line
        line.color.rgb = RGBColor(255, 0, 0)
        
    def slide_f11_emp_tab(self, slide, list_graf):
        f11_emp_tab = self.seguimiento_fs.slides[slide]
        f11_emp_tab.shapes.add_picture(f"{self.path_f11}/{self.fecha_corte}{list_graf[0]}", Cm(4.17), Cm(3.81),  width = Cm(23.41)) 
        f11_emp_tab.shapes.add_picture(f"{self.path_f11}/{self.fecha_corte}{list_graf[1]}", Cm(4.17), Cm(11.48), width = Cm(23.41)) 

    def slide_f11_emp_cli(self, slide, list_graf):
        f11_emp_tab = self.seguimiento_fs.slides[slide]
        f11_emp_tab.shapes.add_picture(f"{self.path_f11}/{self.fecha_corte}{list_graf[0]}", Cm(4.17), Cm(3.81), width = Cm(23.41)) 
        f11_emp_tab.shapes.add_picture(f"{self.path_f11}/{self.fecha_corte}{list_graf[1]}", Cm(4.17), Cm(11.48), width = Cm(23.41)) 

    def slide_f3_costo(self):
        f3 = self.seguimiento_fs.slides[12]
        #Producto
        f3.shapes.add_picture(f"{self.path_f3}/f3_abiertos_fecha_reserva.png", Cm(0.7), Cm(1.8), height=Cm(7.5)) 
        f3.shapes.add_picture(f"{self.path_f3}/f3_tendencia_Producto.png", Cm(13),Cm(1.8), height=Cm(7.5)) 
        f3.shapes.add_picture(f"{self.path_f3}/f3_cerrado_producto_costo.png", Cm(23), Cm(1.8), height=Cm(7.5))
        #Marketplace
        f3.shapes.add_picture(f"{self.path_f3}/f3_abierto_sede.png", Cm(0.7), Cm(10.2), height=Cm(7.5)) 
        f3.shapes.add_picture(f"{self.path_f3}/f3_tendencia_mkp.png", Cm(13), Cm(10.2), height=Cm(7.5)) 
        f3.shapes.add_picture(f"{self.path_f3}/f3_cerrado_mkp_costo.png", Cm(23), Cm(10.2), height=Cm(7.5))

        shapes = f3.shapes
        shape = shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Cm(2.23), Cm(2.98), Cm(7.76), Cm(4.49)
        )
        shape.fill.background()
        line = shape.line
        line.color.rgb = RGBColor(255, 0, 0)



    def slide_f3_tab(self):
        f3 = self.seguimiento_fs.slides[13]
        f3.shapes.add_picture(f"{self.path_f3}/f3_tabla_res_env.png",  Cm(4.58), Cm(5.69), width = Cm(23.41)) 

    def slide_f4(self):
        f4 = self.seguimiento_fs.slides[14]
        f4.shapes.add_picture(f"{self.path_f4}/tendencia_creacion_f4_x_años.png", Cm(2.01),Cm(0.58),width = Cm(16.17), height = Cm(9.08)   )
        f4.shapes.add_picture(f"{self.path_f4}/grafica_f4_sem.png", Cm(19.55), Cm(0.57), width = Cm(13.29), height = Cm(9.6))
        f4.shapes.add_picture(f"{self.path_f4}/grafica_total_por_mes.png", Cm(12.19), Cm(10.07), width = Cm(7.62), height = Cm(7.62)   )
        f4.shapes.add_picture(f"{self.path_f4}/torta.png", Cm(1.34), Cm(10.17), height = Cm(7.13))
        tb_reserva = f4.shapes.add_textbox(Cm(22.59),Cm(14.66),width=Cm(1), height=Cm(0.1))
        tx_text_reserva = tb_reserva.text_frame.add_paragraph()
        print(self.f4_data) 
        tx_text_reserva.text = f"Reservado"# {self.f4_data[3]}"
        tx_text_reserva.font.size= Pt(14)
        
        corchete = f4.shapes.add_textbox(Cm(27.21),Cm(13.4),width=Cm(1), height=Cm(0.1))
        tx_corchete = corchete.text_frame.add_paragraph() 
        tx_corchete.text = "{"
        tx_corchete.font.size= Pt(66)

        estados = f4.shapes.add_textbox(Cm(27.89), Cm(13.91), width=Cm(1), height=Cm(0.1))
        tx_estados = estados.text_frame.add_paragraph() 
        tx_estados.text = f"Tienda self.f4_data[1]\nCD self.f4_data[0]\nDVD self.f4_data[2]\nVenta empresa self.f4_data[3]"
        tx_estados.font.size= Pt(12)

        registrados = f4.shapes.add_textbox(Cm(20.44),Cm(9.97),width=Cm(0.1), height=Cm(0.1))
        tx_registrados = registrados.text_frame.add_paragraph() 
        tx_registrados.text = f"Estado Registrado $21M < 7 días No han pasado a contabilidad"
        tx_registrados.font.size= Pt(13)

        registrados_m = f4.shapes.add_textbox(Cm(20.44),Cm(11.33),width=Cm(0.1), height=Cm(0.1))
        tx_registrados_m = registrados_m.text_frame.add_paragraph() 
        tx_registrados_m.text = f"Estado Registrado $2M >7 días No han pasado a contabilidad,\nya informados a responsables\nUsuarios de creación:\n  Jair Armando Rocha Vargas"
        tx_registrados_m.font.size= Pt(13)       

        shapes = f4.shapes
        left = top = width = height = Cm(1.0)
        shape = shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Cm(14.15), Cm(2.65), Cm(2.76), Cm(1.58)
        )
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)
        line = shape.line
        line.color.rgb = RGBColor(255, 255, 255)


        mes = f4.shapes.add_textbox(Cm(14.04),Cm(1.82),width=Cm(2.07), height=Cm(0.68))
        txt_mes = mes.text_frame.add_paragraph() 
        txt_mes.text = f"Mes Corte: {calendar.month_abbr[self.años[3]]}"
        txt_mes.font.size = Pt(10)
        txt_mes.font.color.rgb = RGBColor(119, 125, 135)

        m_2021 = f4.shapes.add_textbox(Cm(14.04),Cm(2.15),width=Cm(2.07), height=Cm(0.68))
        txt_m_2021 = m_2021.text_frame.add_paragraph() 
        txt_m_2021.text = f"2022: {self.años[0]}"
        txt_m_2021.font.size = Pt(10)
        txt_m_2021.font.color.rgb = RGBColor(0, 0, 0)

        m_2022 = f4.shapes.add_textbox(Cm(14.04),Cm(2.48),width=Cm(2.07), height=Cm(0.68))
        txt_m_2022 = m_2022.text_frame.add_paragraph() 
        txt_m_2022.text = f"2023: {self.años[1]}"
        txt_m_2022.font.size = Pt(10)
        txt_m_2022.font.color.rgb = RGBColor(237, 173, 8)

        dif = f4.shapes.add_textbox(Cm(14.04),Cm(2.81),width=Cm(2.07), height=Cm(0.68))
        txt_dif = dif.text_frame.add_paragraph() 
        txt_dif.text = f"Dif: {self.años[2]}"
        txt_dif.font.size = Pt(10)
        txt_dif.font.color.rgb = RGBColor(0, 179, 80)

    def f4_mot_sede(self):
        f4_mot_sede = self.seguimiento_fs.slides[15]
        f4_mot_sede.shapes.add_picture(f"{self.path_f4}/grafica_total_por_local.png", Cm(0.5),Cm(1.82),width = Cm(31.75), height = Cm(15.88))


    def f4_pos_causa(self):
        f4_pos_causa = self.seguimiento_fs.slides[16]
        f4_pos_causa.shapes.add_picture(f"{self.path_f4}/clasificacion_posibles_causas_22.png", Cm(1.36),Cm(0.38),width = Cm(31.91), height = Cm(17.32))
        titulo = f4_pos_causa.shapes.add_textbox(Cm(0.4),Cm(-0.79),width=Cm(0.1), height=Cm(0.1))
        tx_titulo = titulo.text_frame.add_paragraph()      

    def f4_mes_moti(self):
        f4_mes_moti = self.seguimiento_fs.slides[17]
        f4_mes_moti.shapes.add_picture(f"{self.path_f4}/mes_f4_motivo_compañia.png", Cm(0.34),Cm(2.88), height = Cm(14.02))
        f4_mes_moti.shapes.add_picture(f"{self.path_f4}/cd_mm.png", Cm(17.82),Cm(0.48), height = Cm(8.27))
        f4_mes_moti.shapes.add_picture(f"{self.path_f4}/tienda_mes_motivo.png", Cm(17.82),Cm(9.27), height = Cm(8.38))

    def f4_linea_año(self):
        f4_linea_año = self.seguimiento_fs.slides[18]
        f4_linea_año.shapes.add_picture(f"{self.path_f4}/torta_linea.png", Cm(2.06),Cm(2.4),width = Cm(12.73), height = Cm(14.55))
        f4_linea_año.shapes.add_picture(f"{self.path_f4}/grafica_linea_x_mes.png", Cm(14.79),Cm(1.89),width = Cm(17.46), height = Cm(15.27))

    def f4_linea_motivo_mes(self):
        f4_linea_motivo_mes = self.seguimiento_fs.slides[19]
        f4_linea_motivo_mes.shapes.add_picture(f"{self.path_f4}/fig_linea_mes_mot.png", Cm(8.73), Cm(1.19), height = Cm(16.4))


    def f4_linea_motivo(self):
        f4_linea_motivo = self.seguimiento_fs.slides[20]
        f4_linea_motivo.shapes.add_picture(f"{self.path_f4}/f4_linea_motivo.png", Cm(0),Cm(2.66), height = Cm(13.74))
        f4_linea_motivo.shapes.add_picture(f"{self.path_f4}/linea_local.png", Cm(17),Cm(2.66), height = Cm(13.74))

    def f4_averias(self):
        f4_linea_motivo = self.seguimiento_fs.slides[21]
        f4_linea_motivo.shapes.add_picture(f"{self.path_f4}/grafica_averia_x_mes_y_marca.png", Cm(0),Cm(2.54),width = Cm(16.64), height = Cm(13.82))
        f4_linea_motivo.shapes.add_picture(f"{self.path_f4}/tabla_averias.png", Cm(16.02),Cm(2.6),width = Cm(17.44))
        #TODO agregar graficas, marcas x locales

    def f4_calidad(self):
        f4_calidad = self.seguimiento_fs.slides[22]
        f4_calidad.shapes.add_picture(f"{self.path_f4}/marca_calidad.png", Cm(0),Cm(2.4),width = Cm(16.64), height = Cm(13.82))
        f4_calidad.shapes.add_picture(f"{self.path_f4}/tabla_calidad.png", Cm(16.64),Cm(6.11),width = Cm(16.64), height = Cm(6.84))

    def f4_panatallas_rotas(self):
        f4_panatallas_rotas = self.seguimiento_fs.slides[23]
        f4_panatallas_rotas.shapes.add_picture(f"{self.path_f4}/pantallas_rotas.png", Cm(0.5),Cm(1.2), height = Cm(16.25))
        f4_panatallas_rotas.shapes.add_picture(f"{self.path_f4}/tabla_pantallas_rotas.png", Cm(13.7), Cm(1.2),width = Cm(21.23))
        f4_panatallas_rotas.shapes.add_picture(f"{self.path_f4}/tabla_prod_pant.png", Cm(13.7), Cm(6.1), width = Cm(7.99))
        f4_panatallas_rotas.shapes.add_picture(f"{self.path_f4}/tabla_loc_pant.png", Cm(22.26),Cm(7.2), width = Cm(11.45))

    def f4_anulados(self):
        f4_anulados = self.seguimiento_fs.slides[24]
        f4_anulados.shapes.add_picture(f"{self.path_f4}/pie_anulados.png", Cm(3.48),Cm(1.99), height = Cm(7.46))
        f4_anulados.shapes.add_picture(f"{self.path_f4}/anulados_mes.png", Cm(13.33),Cm(1.86),height = Cm(7.46))
        f4_anulados.shapes.add_picture(f"{self.path_f4}/tabla_anulados.png", Cm(3.48),Cm(11.04),width = Cm(26.9))

    def f4_enviados(self):
        f4_enviados = self.seguimiento_fs.slides[25]
        f4_enviados.shapes.add_picture(f"{self.path_f4}/tabla_enviados.png", Cm(3.48),Cm(11.04), width = Cm(26.9))
        f4_enviados.shapes.add_picture(f"{self.path_f4}/enviados.png",  Cm(6.35), Cm(1.77), height = Cm(7.94))

    def f4_registrados(self):
        print('Activar slide de registrados')
        f4_registrados = self.seguimiento_fs.slides[26]
        # f4_registrados.shapes.add_picture(f"{self.path_f4}/tabla_registrados.png", Cm(3.48), Cm(11.04) ,width = Cm(26.9), height = Cm(3.97))
        # f4_registrados.shapes.add_picture(f"{self.path_f4}/registrados.png", Cm(7.03), Cm(2.24), width = Cm(26.9))

    def f4_reservados(self):
        f4_reservados = self.seguimiento_fs.slides[27]
        f4_reservados.shapes.add_picture(f"{self.path_f4}/tabla_reservados.png",  Cm(3.48),Cm(11.04), width = Cm(26.9))
    
    def f4_detalle(self):
        f4_detalle = self.seguimiento_fs.slides[28]
        f4_detalle.shapes.add_picture(f"{self.path_f4}/tabla_materiales.png", Cm(1.32),Cm(3.14),width = Cm(13.23), height = Cm(3.84))
        f4_detalle.shapes.add_picture(f"{self.path_f4}/tabla_act_fijo.png", Cm(1.32),Cm(8.73),width = Cm(13.23), height = Cm(3.44))
        f4_detalle.shapes.add_picture(f"{self.path_f4}/tabla_receta.png", Cm(1.32),Cm(13.79),width = Cm(13.23), height = Cm(2.38))
        f4_detalle.shapes.add_picture(f"{self.path_f4}/tabla_gasto_int.png", Cm(18.43),Cm(4.33),width = Cm(13.23), height = Cm(3.84))
        f4_detalle.shapes.add_picture(f"{self.path_f4}/tabla_bolsa.png", Cm(18.43),Cm(11.02),width = Cm(13.23), height = Cm(3.31))

    def get_all_f4_slides(self):
        self.slide_f4()
        self.f4_pos_causa()
        self.f4_mot_sede()
        self.f4_mes_moti()
        self.f4_linea_año()
        self.f4_linea_motivo()
        self.f4_linea_motivo_mes()
        self.f4_averias()
        self.f4_calidad()
        self.f4_panatallas_rotas()
        self.f4_anulados()
        self.f4_enviados()
        self.f4_registrados()
        self.f4_reservados()
        self.f4_detalle()

    def get_all_f3_slides(self):
        self.slide_f3_costo()
        self.slide_f3_tab()

    def get_all_f11_slides(self):
        self.slide_f11_empresa(2, var_f11["graf_monto"])
        self.slide_f11_empresa(3, var_f11["graf_cant"])
        self.slide_f11_emp_tab(5, var_f11["tab_emp"])
        self.slide_f11_emp_tab(6, var_f11["tab_emp_loc"])
        self.slide_f11_emp_cli(7, var_f11["tab_cl"])
        self.slide_f11_emp_cli(8, var_f11["tab_cl_loc"])

    def save_ppt(self):
        path_ppt = f"{var_global['path_ppts']}/seguimiento_fs.pptx"
        self.seguimiento_fs.save(path_ppt)
        return path_ppt
    